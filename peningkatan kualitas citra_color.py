import cv2
import numpy as np
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Inches, Pt

# === Folder output ===
output_dir = "outputs_pertemuan5_color"
os.makedirs(output_dir, exist_ok=True)

# === Load gambar berwarna ===
SOURCE_IMAGE_PATH = "foto_saya.jpg"   # ganti dengan gambar kamu
src = cv2.imread(SOURCE_IMAGE_PATH)   # BGR
src_rgb = cv2.cvtColor(src, cv2.COLOR_BGR2RGB)  # biar tampil normal

# === Fungsi simpan histogram RGB ===
def plot_hist_rgb(img, path, title):
    color = ('r','g','b')
    plt.figure(figsize=(6,4))
    for i,col in enumerate(color):
        hist = cv2.calcHist([img],[i],None,[256],[0,256])
        plt.plot(hist,color=col)
        plt.xlim([0,256])
    plt.title(title)
    plt.xlabel("Intensitas")
    plt.ylabel("Jumlah Piksel")
    plt.tight_layout()
    plt.savefig(path)
    plt.close()

# === Equalisasi Histogram (pada channel Y) ===
def equalize_color(img):
    img_ycrcb = cv2.cvtColor(img, cv2.COLOR_RGB2YCrCb)
    y,cr,cb = cv2.split(img_ycrcb)
    y_eq = cv2.equalizeHist(y)
    img_ycrcb_eq = cv2.merge([y_eq, cr, cb])
    img_eq = cv2.cvtColor(img_ycrcb_eq, cv2.COLOR_YCrCb2RGB)
    return img_eq

# === Histogram Matching per channel ===
def match_histograms_color(source, template):
    matched = np.zeros_like(source)
    for i in range(3): # R,G,B
        s_vals, bin_idx, s_counts = np.unique(source[:,:,i].ravel(), return_inverse=True, return_counts=True)
        t_vals, t_counts = np.unique(template[:,:,i].ravel(), return_counts=True)

        s_quantiles = np.cumsum(s_counts).astype(np.float64)
        s_quantiles /= s_quantiles[-1]
        t_quantiles = np.cumsum(t_counts).astype(np.float64)
        t_quantiles /= t_quantiles[-1]

        interp_t = np.interp(s_quantiles, t_quantiles, t_vals)
        matched[:,:,i] = interp_t[bin_idx].reshape(source[:,:,i].shape).astype(np.uint8)
    return matched

# === Buat citra target lebih terang/kontras ===
def make_target_color(src):
    lab = cv2.cvtColor(src, cv2.COLOR_RGB2LAB)
    l,a,b = cv2.split(lab)
    l = cv2.equalizeHist(l)
    lab_eq = cv2.merge([l,a,b])
    return cv2.cvtColor(lab_eq, cv2.COLOR_LAB2RGB)

# Proses
target = make_target_color(src_rgb)
src_eq = equalize_color(src_rgb)
src_matched = match_histograms_color(src_rgb, target)

# Simpan hasil + histogram
def save_results(img, hist_title, prefix):
    img_path = os.path.join(output_dir, f"{prefix}.png")
    hist_path = os.path.join(output_dir, f"{prefix}_hist.png")
    cv2.imwrite(img_path, cv2.cvtColor(img, cv2.COLOR_RGB2BGR))
    plot_hist_rgb(img, hist_path, hist_title)
    return img_path, hist_path

orig_img, orig_hist = save_results(src_rgb, "Histogram Citra Asli", "original")
equal_img, equal_hist = save_results(src_eq, "Histogram Hasil Equalisasi", "equalized")
target_img, target_hist = save_results(target, "Histogram Citra Target", "target")
matched_img, matched_hist = save_results(src_matched, "Histogram Hasil Matching", "matched")

# === Buat PowerPoint ===
prs = Presentation()

# Slide Judul
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Tugas Mandiri Pertemuan ke-5 (Versi Warna)"
slide.placeholders[1].text = "Nama: [Isi Nama Anda]\nNIM: [Isi NIM Anda]"

def add_slide(title_text, img_path, hist_path, note_text):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tx = s.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.6))
    tx.text_frame.text = title_text
    s.shapes.add_picture(img_path, Inches(0.4), Inches(0.9), width=Inches(4.5))
    s.shapes.add_picture(hist_path, Inches(5.5), Inches(0.9), width=Inches(4.0))
    tb = s.shapes.add_textbox(Inches(0.4), Inches(4.6), Inches(9), Inches(1.5))
    tb.text_frame.text = note_text

add_slide("Citra Asli + Histogram", orig_img, orig_hist, "Histogram menunjukkan distribusi warna asli.")
add_slide("Equalisasi Histogram", equal_img, equal_hist, "Kontras meningkat dengan tetap menjaga warna asli.")
add_slide("Citra Target + Histogram", target_img, target_hist, "Target digunakan sebagai acuan tampilan.")
add_slide("Spesifikasi Histogram", matched_img, matched_hist, "Histogram dicocokkan agar mirip target.")

# Slide Kesimpulan
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Kesimpulan"
s.placeholders[1].text = ("1. Histogram RGB menunjukkan distribusi tiap warna.\n"
                          "2. Equalisasi pada channel luminance meningkatkan kontras tanpa merusak warna.\n"
                          "3. Spesifikasi histogram menghasilkan tampilan citra mirip target.")

pptx_path = os.path.join(output_dir, "hasil_tugas_pertemuan5_warna.pptx")
prs.save(pptx_path)

print("✅ Selesai! File PowerPoint tersimpan di:", pptx_path)
