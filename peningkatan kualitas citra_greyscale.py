# === Import library ===
from PIL import Image
import numpy as np
import matplotlib.pyplot as plt
import os

# untuk buat pptx
from pptx import Presentation
from pptx.util import Inches, Pt

# === Folder output ===
output_dir = "outputs_pertemuan5"
os.makedirs(output_dir, exist_ok=True)

# === Fungsi load citra ===
def open_image_as_gray(path):
    im = Image.open(path).convert("L")  # ubah ke grayscale
    return np.array(im)

# >>> Ganti path di sini dengan gambar Anda sendiri <<<
SOURCE_IMAGE_PATH = "foto_saya.jpg"   # misalnya foto landscape Anda
src = open_image_as_gray(SOURCE_IMAGE_PATH)

# === Fungsi simpan citra & histogram ===
def save_image(arr, path):
    Image.fromarray(arr).save(path)

def plot_and_save_histogram(arr, path, title):
    plt.figure(figsize=(6,3.5))
    plt.hist(arr.ravel(), bins=256, range=(0,255))
    plt.title(title)
    plt.xlabel("Intensitas")
    plt.ylabel("Jumlah Piksel")
    plt.tight_layout()
    plt.savefig(path)
    plt.close()

# === Fungsi equalization ===
def equalize_histogram_manual(arr):
    hist, bins = np.histogram(arr.flatten(), 256, [0,256])
    cdf = hist.cumsum()
    cdf_masked = np.ma.masked_equal(cdf, 0)
    cdf_masked = (cdf_masked - cdf_masked.min()) * 255 / (cdf_masked.max() - cdf_masked.min())
    cdf_final = np.ma.filled(cdf_masked, 0).astype('uint8')
    img_eq = cdf_final[arr]
    return img_eq

# === Fungsi spesifikasi histogram ===
def match_histograms_manual(source, template):
    s_vals, bin_idx, s_counts = np.unique(source.ravel(), return_inverse=True, return_counts=True)
    t_vals, t_counts = np.unique(template.ravel(), return_counts=True)

    s_quantiles = np.cumsum(s_counts).astype(np.float64)
    s_quantiles /= s_quantiles[-1]
    t_quantiles = np.cumsum(t_counts).astype(np.float64)
    t_quantiles /= t_quantiles[-1]

    interp_t = np.interp(s_quantiles, t_quantiles, t_vals)
    matched = interp_t[bin_idx].reshape(source.shape).astype(np.uint8)
    return matched

# === Buat citra target sederhana (lebih terang / kontras) ===
def make_target_from(src):
    arr = src.astype(np.float32)
    p2, p98 = np.percentile(arr, (2, 98))
    arr = (arr - p2) * 255.0 / (p98 - p2 + 1e-6)
    arr = 255.0 * (np.clip(arr, 0, 255) / 255.0) ** 0.8
    return np.clip(arr, 0, 255).astype(np.uint8)

target = make_target_from(src)

# === Proses gambar ===
src_eq = equalize_histogram_manual(src)
src_matched = match_histograms_manual(src, target)

# === Simpan semua hasil ===
files = {}
files['orig_img'] = os.path.join(output_dir, "original.png")
files['orig_hist'] = os.path.join(output_dir, "original_hist.png")
files['equal_img'] = os.path.join(output_dir, "equalized.png")
files['equal_hist'] = os.path.join(output_dir, "equalized_hist.png")
files['target_img'] = os.path.join(output_dir, "target.png")
files['target_hist'] = os.path.join(output_dir, "target_hist.png")
files['matched_img'] = os.path.join(output_dir, "matched.png")
files['matched_hist'] = os.path.join(output_dir, "matched_hist.png")

save_image(src, files['orig_img'])
plot_and_save_histogram(src, files['orig_hist'], "Histogram Citra Asli")

save_image(src_eq, files['equal_img'])
plot_and_save_histogram(src_eq, files['equal_hist'], "Histogram Hasil Equalisasi")

save_image(target, files['target_img'])
plot_and_save_histogram(target, files['target_hist'], "Histogram Citra Target")

save_image(src_matched, files['matched_img'])
plot_and_save_histogram(src_matched, files['matched_hist'], "Histogram Hasil Matching")

# === Buat PowerPoint ===
prs = Presentation()

# Slide Judul
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Tugas Mandiri Pertemuan ke-5\nAnalisis Histogram, Equalisasi, dan Spesifikasi"
slide.placeholders[1].text = "Nama: [Isi Nama Anda]\nNIM: [Isi NIM Anda]"

# Fungsi helper: tambah slide gambar + histogram
def add_image_hist_slide(title_text, img_path, hist_path, note_text):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    tx = s.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.6))
    tf = tx.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.size = Pt(18)

    s.shapes.add_picture(img_path, Inches(0.4), Inches(0.9), width=Inches(4.5))
    s.shapes.add_picture(hist_path, Inches(5.5), Inches(0.9), width=Inches(4.0))

    tb = s.shapes.add_textbox(Inches(0.4), Inches(4.6), Inches(9), Inches(1.5))
    tb.text = note_text

# Tambah slide untuk setiap tahap
add_image_hist_slide("Citra Asli dan Histogram", files['orig_img'], files['orig_hist'],
                     "Histogram menunjukkan distribusi intensitas asli.")
add_image_hist_slide("Hasil Equalisasi Histogram", files['equal_img'], files['equal_hist'],
                     "Equalisasi membuat distribusi lebih merata sehingga kontras meningkat.")
add_image_hist_slide("Citra Target dan Histogram", files['target_img'], files['target_hist'],
                     "Citra target digunakan sebagai acuan kontras/kecerahan.")
add_image_hist_slide("Hasil Spesifikasi Histogram", files['matched_img'], files['matched_hist'],
                     "Histogram citra asli disesuaikan agar mirip dengan citra target.")

# Slide Kesimpulan
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Kesimpulan"
s.placeholders[1].text = ("1. Histogram membantu analisis tingkat terang/gelap citra.\n"
                          "2. Equalisasi meningkatkan kontras dengan meratakan distribusi.\n"
                          "3. Spesifikasi membuat citra menyerupai target yang diinginkan.")

# Simpan PPTX
pptx_path = os.path.join(output_dir, "hasil_tugas_pertemuan5.pptx")
prs.save(pptx_path)

print("=== Proses selesai ===")
print("File PowerPoint tersimpan di:", pptx_path)
